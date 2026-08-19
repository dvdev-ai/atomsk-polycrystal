#!/usr/bin/env python3
"""Генерация презентации PowerPoint для защиты курсовой."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
FIG = ROOT / "figures"
OUT = ROOT / "presentation"
OUT.mkdir(exist_ok=True)

# Цветовая схема
NAVY = RGBColor(0x1A, 0x36, 0x5D)
TEAL = RGBColor(0x00, 0x96, 0x88)
LIGHT = RGBColor(0xF0, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)


def set_title_style(shape, size: int = 32, color=NAVY) -> None:
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.5), Inches(2))
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(8.5), Inches(1.2))
    st = sub.text_frame
    st.text = subtitle
    sp = st.paragraphs[0]
    sp.font.size = Pt(18)
    sp.font.color.rgb = TEAL


def add_section(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, 0, Inches(2.8), prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.95), Inches(9), Inches(0.8))
    tb.text_frame.text = title
    p = tb.text_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_image_slide(prs: Presentation, title: str, image: Path, caption: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    hdr = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    hdr.text_frame.text = title
    set_title_style(hdr, 22)
    if image.exists():
        slide.shapes.add_picture(str(image), Inches(0.5), Inches(1.0), width=Inches(9))
    if caption:
        cap = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(9), Inches(0.5))
        cap.text_frame.text = caption
        p = cap.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY


def add_bullets(prs: Presentation, title: str, bullets: list[str], image: Path | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # заголовок
    hdr = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.7))
    hdr.text_frame.text = title
    set_title_style(hdr, 24)

    left = Inches(0.55)
    width = Inches(5.2) if image and image.exists() else Inches(9)
    body = slide.shapes.add_textbox(left, Inches(1.15), width, Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.space_after = Pt(8)

    if image and image.exists():
        slide.shapes.add_picture(str(image), Inches(5.9), Inches(1.2), width=Inches(3.8))


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Программный комплекс ATOMSK",
        "Поликристаллические Cu и Al при одноосном растяжении\n"
        "LAMMPS · OVITO · Диаграммы σ–ε (0, 300, 600 K)",
    )

    add_bullets(
        prs,
        "Цель работы",
        [
            "Изучить назначение и возможности ATOMSK в атомистическом моделировании.",
            "Построить модели поликристаллических Cu и Al (метод Вороного).",
            "Провести MD-расчёт одноосного растяжения в LAMMPS.",
            "Получить диаграммы «напряжение–деформация» при 0, 300 и 600 K.",
            "Визуализировать структуру и деформирование в OVITO.",
        ],
    )

    add_section(prs, "ATOMSK: назначение и возможности")

    add_bullets(
        prs,
        "Что такое ATOMSK?",
        [
            "Инструмент командной строки для создания и преобразования атомных конфигураций.",
            "Форматы: LAMMPS, CFG, XSF, POSCAR, PDB и др.",
            "Режим --polycrystal: поликристаллы методом тесселяции Вороного.",
            "Свойство grainID — номер зерна (удобно для OVITO).",
            "Сайт: atomsk.univ-lille.fr (P. Hirel, Univ. Lille).",
        ],
    )

    add_bullets(
        prs,
        "Альтернативные подходы",
        [
            "Ручная сборка в VESTA / Crystal Maker — трудоёмко для поликристаллов.",
            "Пакеты Neper, Dream.3D — зерновая структура + конвертация в MD.",
            "Встроенное create_atoms в LAMMPS — только однородные/регулярные структуры.",
            "ATOMSK — быстрый способ получить реалистичный поликристалл из seed-ячейки.",
        ],
    )

    add_section(prs, "Методика расчёта")

    add_bullets(
        prs,
        "Построение модели (ATOMSK)",
        [
            "Cu: FCC, a = 3.615 Å; Al: FCC, a = 4.046 Å.",
            "Куб 60×60×60 Å, 8 зёрен (случайные ориентации).",
            "Команда: atomsk --polycrystal seed.xsf params.txt out.cfg -wrap",
            "Экспорт: atomsk out.cfg lammps Cu_poly.data",
            "~18 000 атомов (Cu), ~13 000 (Al).",
        ],
    )

    add_bullets(
        prs,
        "Расчёт в LAMMPS",
        [
            "Потенциалы EAM: Cu_mishin1.eam.alloy, Al_zhou.eam.alloy.",
            "Минимизация энергии → NPT-релаксация при заданной T.",
            "Одноосное растяжение: fix deform x, боковые грани NPT (pyy, pzz → 0).",
            "Скорость деформации ε̇ ≈ 10⁹ s⁻¹ (типично для MD).",
            "Выход: stress_strain.dat (ε, σ_xx в ГПа).",
        ],
    )

    add_section(prs, "Визуализация OVITO")

    ovito_cu = FIG / "ovito_Cu_polycrystal.png"
    ovito_al = FIG / "ovito_Al_polycrystal.png"
    ovito_tens = FIG / "ovito_tension_Cu_T300K.png"
    if not ovito_tens.exists():
        ovito_tens = FIG / "ovito_tension_quick_Cu_300K.png"
    if not ovito_tens.exists():
        ovito_tens = next(iter(sorted(FIG.glob("ovito_tension_*.png"))), Path())

    if ovito_cu.exists():
        add_image_slide(
            prs,
            "Поликристалл Cu (раскраска по grainID)",
            ovito_cu,
            "ATOMSK → CFG → OVITO, 8 зёрен, 60×60×60 Å",
        )
    if ovito_al.exists():
        add_image_slide(
            prs,
            "Поликристалл Al (раскраска по grainID)",
            ovito_al,
            "Метод Вороного, случайные ориентации зёрен",
        )
    if ovito_tens.exists():
        add_image_slide(
            prs,
            "Деформированное состояние (LAMMPS → OVITO)",
            ovito_tens,
            "Cu, 300 K: финальный кадр dump_tension.lammpstrj, раскраска по координате X (направление растяжения)",
        )

    add_section(prs, "Результаты")

    img_cu = FIG / "stress_strain_Cu.png"
    img_al = FIG / "stress_strain_Al.png"
    img_all = FIG / "stress_strain_combined.png"
    img_sum = FIG / "summary_table.png"

    if img_all.exists():
        add_image_slide(
            prs,
            "Диаграммы σ–ε: Cu и Al (0, 300, 600 K)",
            img_all,
            "Поликристалл, ε̇ = 10⁹ s⁻¹, ε до ~3%, потенциалы EAM",
        )
    if img_cu.exists():
        add_image_slide(
            prs,
            "Медь (Cu): напряжение–деформация",
            img_cu,
            "σ_max снижается с ростом T (0 → 600 K)",
        )
    if img_al.exists():
        add_image_slide(
            prs,
            "Алюминий (Al): напряжение–деформация",
            img_al,
            "Три температуры, одноосное растяжение по X",
        )
    if img_sum.exists():
        add_image_slide(prs, "Сводка численных результатов", img_sum, "")

    add_bullets(
        prs,
        "Выводы",
        [
            "ATOMSK эффективен для генерации поликристаллов без ручной расстановки атомов.",
            "Связка ATOMSK → LAMMPS → OVITO — стандартный пайплайн в материаловедении.",
            "Получены кривые σ–ε для Cu и Al при трёх температурах.",
            "Для количественного согласия с экспериментом нужны большие образцы и медленнее ε̇.",
        ],
    )

    add_bullets(
        prs,
        "Спасибо за внимание!",
        [
            "Вопросы?",
            "",
            "Материалы проекта: ~/Downloads/atomsk-polycrystal-tension",
            "README_RU.md — пошаговая инструкция",
        ],
    )

    path = OUT / "ATOMSK_курсовая.pptx"
    prs.save(path)
    # Отдельная папка для просмотра
    view_dir = ROOT / "Презентация"
    view_dir.mkdir(exist_ok=True)
    view_path = view_dir / "ATOMSK_курсовая.pptx"
    import shutil

    shutil.copy2(path, view_path)
    shutil.copy2(path, Path.home() / "Downloads" / "ATOMSK_курсовая.pptx")
    print(f"Презентация сохранена: {path}")
    print(f"Для просмотра: {view_path}")


if __name__ == "__main__":
    main()
